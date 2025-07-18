import os

def generate_ppt(input_filename="input.xlsx"):
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches,Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    input_path = os.path.join(os.path.dirname(__file__), input_filename)
    output_path = os.path.join(os.path.dirname(__file__), "ppt.pptx")

    df = pd.read_excel(input_path)

    TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template", "template.pptx")

    prs = Presentation(TEMPLATE_PATH)
    
        #slide 1 to 8 are constant

    """# Slide 9"""

    df = pd.read_excel(input_path , sheet_name="Slide 9 - Model Portfolio Per", header=None)

    # Extract header (row 1, columns 1 to 5)
    columns = df.iloc[1, 1:6]

    # Extract data (rows 2 onwards, same cols)
    table = df.iloc[2:, 1:6]
    table.columns = columns
    table = table.reset_index(drop=True)

    # Assuming table1 is already assigned:
    slide = prs.slides[8]
    table_shape = slide.shapes[1]
    table1 = table_shape.table

    # Fill table body (starting from row 1)
    for i in range(len(table)):
        if i + 1 >= len(table1.rows):
            continue  # Skip if table doesn't have enough rows

        for j in range(len(table.columns)):
            if j >= len(table1.columns):
                continue  # Skip if table doesn't have enough columns

            val = table.iloc[i, j]
            cell = table1.cell(i + 1, j)

            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()

            # Format value
            if j == 0:
                run.text = str(val)
            else:
                try:
                    if i >= len(table) - 2:
                        run.text = f"{float(val)*100:.2f}%"
                    else:
                        run.text = f"{float(val)*100:.1f}%"
                except:
                    run.text = str(val)

            run.font.name = 'Calibri'
            run.font.size = Pt(10.5)

            # Make bottom 2 rows bold
            if i >= len(table) - 2:
                run.font.bold = True
            else:
                run.font.bold = False

            # Left-most column → white bold
            if j == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)  # white
            else:
                try:
                    if float(val) < 0:
                        run.font.color.rgb = RGBColor(255, 0, 0)  # red
                    else:
                        run.font.color.rgb = RGBColor(0, 0, 0)    # black
                except:
                    run.font.color.rgb = RGBColor(0, 0, 0)

            p.alignment = PP_ALIGN.CENTER

    # Fill header row (row 0)
    for j in range(len(table.columns)):
        if j >= len(table1.columns):
            continue

        cell = table1.cell(0, j)
        p = cell.text_frame.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = str(table.columns[j])
        run.font.name = 'Calibri'
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)  # black
        p.alignment = PP_ALIGN.CENTER

    """#Slide 10"""

    # Load the specific sheet
    df = pd.read_excel(input_path , sheet_name="Slide 10 - Current Port+ Expens", header=None)

    # Extract the table: rows from 3 onwards, columns B to E (i.e., 1 to 4)
    table = df.iloc[3:10, 1:7]

    # Set the correct headers (2nd row, index 2)
    table.columns = df.iloc[2, 1:7]

    # Reset index if needed
    table = table.reset_index(drop=True)


    slide = prs.slides[9]
    table_shape = slide.shapes[1]
    table1 = table_shape.table

    # Identify column indexes for formatting
    weights_col = table.columns.get_loc('Weights')
    expense_col = table.columns.get_loc('Direct (Expense Ratio)')

    # Body rows
    for i, row in enumerate(table.itertuples(index=False)):
        for j, val in enumerate(row):
            # Format weights/expense cols as percentage
            if j == weights_col or j == expense_col:
                val = f"{float(val)*100:.2f}%" if pd.notna(val) and str(val).replace('.', '', 1).isdigit() else val

            cell = table1.cell(i+1, j)
            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(val) if pd.notna(val) else ""
            run.font.name = 'Calibri'
            run.font.size = Pt(9.2)
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False
            p.alignment = PP_ALIGN.CENTER

    # Header row
    for j, val in enumerate(table.columns):
        cell = table1.cell(0, j)
        p = cell.text_frame.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = str(val)
        run.font.name = 'Calibri'
        run.font.size = Pt(9.5)
        run.font.bold = True
        run.font.italic = True
        run.font.underline = True
        run.font.color.rgb = RGBColor(255, 255, 255)  # White
        p.alignment = PP_ALIGN.CENTER

    """# Slide 11"""

    df = pd.read_excel(input_path , sheet_name="Slide 11- Key Highlights", header=None)
    # Extract sector data: rows 2 to 6 and 8 to 10, columns 1 and 2
    sector_top = df.iloc[2:7, [1, 2]]
    sector_bottom = df.iloc[8:11, [1, 2]]

    # Concatenate top and bottom parts
    sector_table = pd.concat([sector_top, sector_bottom]).reset_index(drop=True)

    # Set headers manually (row 1)
    sector_table.columns = ['Sector', 'Weightage (%)']

    slide = prs.slides[10]  # Slide 11
    table_shape = slide.shapes[2]  # Table 11
    table1 = table_shape.table

    # Format the weight column as percentage with 1 decimal
    sector_table['Weightage (%)'] = sector_table['Weightage (%)'].apply(
        lambda x: f"{round(x * 100, 1)}%" if pd.notna(x) else ""
    )

    # Push values into last 2 columns (Top 5 Sectors)
    for i in range(5):  # Only 5 visible rows in table
        for j in range(2):  # Sector name + weight
            val = sector_table.iloc[i, j]
            cell = table1.cell(i + 1, j + 4)  # Offset by 1 for header row, +4 for last 2 cols

            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(val)
            run.font.name = 'Calibri'
            run.font.size = Pt(10.5)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

    """#Slide 12"""

    df = pd.read_excel(input_path , sheet_name="Slide 12 - Performance", header=None)
    table1 = df.iloc[1:5, 1:6]
    table1.columns = table1.iloc[0]  # set proper column headers
    table1 = table1.drop(index=1)    # drop redundant header row if any
    table1 = table1.reset_index(drop=True)

    table2 = df.iloc[8:14, 1:4]
    table2.columns = df.iloc[7, 1:4]  # row 7 is the header
    table2 = table2.reset_index(drop=True)

    table1_shape = prs.slides[11].shapes[2]
    ppt_table1 = table1_shape.table

    for i in range(len(table1)):
        for j in range(len(table1.columns)):
            val = table1.iloc[i, j]
            cell = ppt_table1.cell(i + 1, j)

            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            try:
                val_f = float(val) * 100
                val_fmt = f"{val_f:.2f}%"
            except:
                val_fmt = str(val)
            run.text = val_fmt
            run.font.name = 'Calibri'
            run.font.size = Pt(10.5)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

    # Headers
    for j, col in enumerate(table1.columns):
        cell = ppt_table1.cell(0, j)
        p = cell.text_frame.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = str(col)
        run.font.name = 'Calibri'
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    table2_shape = prs.slides[11].shapes[3]
    ppt_table2 = table2_shape.table

    for i in range(len(table2)):
        for j in range(len(table2.columns)):
            val = table2.iloc[i, j]
            cell = ppt_table2.cell(i + 1, j)

            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            try:
                if j != 0:  # convert weight & performance to percentage
                    val_f = float(val) * 100
                    val_fmt = f"{val_f:.2f}%"
                else:
                    val_fmt = str(val)
            except:
                val_fmt = str(val)

            run.text = val_fmt
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

    # Header row
    for j, col in enumerate(table2.columns):
        cell = ppt_table2.cell(0, j)
        p = cell.text_frame.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = str(col)
        run.font.name = 'Calibri'
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    """# Slide 13"""

    df = pd.read_excel(input_path , sheet_name="Slide 13 - Canara Emerging", header=None)
    slide = prs.slides[12]

    # 1. Key Details
    key_details_heading = df.iloc[3, 1:3].tolist()
    key_details = df.iloc[4:10, 1:3].dropna().reset_index(drop=True)
    key_details.columns = key_details_heading

    # Format Inception date if it's a full datetime
    if 'Inception' in key_details[key_details.columns[0]].values:
        idx = key_details[key_details[key_details.columns[0]] == 'Inception'].index[0]
        val = key_details.loc[idx, key_details.columns[1]]
        try:
            formatted = pd.to_datetime(str(val), errors='coerce').strftime('%b-%y')
            key_details.loc[idx, key_details.columns[1]] = formatted
        except:
            pass

    # 2. Portfolio Aggregates
    portfolio_aggregates_heading = df.iloc[3, 5:7].tolist()
    portfolio_aggregates = df.iloc[4:11, 5:7].dropna().reset_index(drop=True)
    portfolio_aggregates.columns = portfolio_aggregates_heading
    portfolio_aggregates.iloc[3:, 1] = pd.to_numeric(portfolio_aggregates.iloc[3:, 1], errors='coerce').apply(lambda x: f"{round(x * 100)}%" if pd.notna(x) else "-")

    # 3. Portfolio Concentration
    concentration_heading = df.iloc[3, [9, 11]].tolist()
    concentration_heading = ['Metric', 'Value'] if any(pd.isna(concentration_heading)) else concentration_heading
    concentration = df.iloc[4:9, [9, 11]].dropna().reset_index(drop=True)
    concentration.columns = concentration_heading
    concentration_vals = pd.to_numeric(concentration.iloc[:, 1], errors='coerce')
    concentration.iloc[:, 1] = [
        f"{round(val * 100, 1)}%" if pd.notna(val) and val < 10 else
        f"{round(val, 1)}x" if pd.notna(val) else "-"
        for val in concentration_vals
    ]

    # 4. Top Sectors
    top_sectors_heading = df.iloc[3, 13:15].tolist()
    top_sectors_heading = ['Sector', 'Weight'] if any(pd.isna(top_sectors_heading)) else top_sectors_heading
    top_sectors = df.iloc[4:10, 13:15].dropna().reset_index(drop=True)
    top_sectors.columns = top_sectors_heading
    top_sectors.iloc[:, 1] = pd.to_numeric(top_sectors.iloc[:, 1], errors='coerce').apply(lambda x: f"{round(x * 100, 1)}%" if pd.notna(x) else "-")

    # 5. CY Returns (Corrected)
    cy_returns_heading = df.iloc[17, 1:5].tolist()  # B to E heading
    cy_returns = df.iloc[18:27, 1:5].dropna().reset_index(drop=True)  # B to E data (Year + 3 metrics)
    cy_returns.columns = cy_returns_heading
    for col in cy_returns.columns[1:]:  # Skip 'Year'
        cy_returns[col] = pd.to_numeric(cy_returns[col], errors='coerce').apply(lambda x: f"{round(x * 100, 1)}%" if pd.notna(x) else "-")

    # 6. Top Holdings
    top_holdings_heading = df.iloc[16, [7, 8]].tolist()
    top_holdings_heading = ['Top Holdings', '% Allocation'] if any(pd.isna(top_holdings_heading)) else top_holdings_heading
    top_holdings = df.iloc[17:25, [7, 8]].dropna().reset_index(drop=True)
    top_holdings.columns = top_holdings_heading
    top_holdings.iloc[:, 1] = pd.to_numeric(top_holdings.iloc[:, 1], errors='coerce').apply(lambda x: round(x, 2) if pd.notna(x) else "-")

    # Slide Heading
    slide_heading = df.iloc[1, 1]
    slide = prs.slides[12]

    title_shape = slide.shapes[0]
    title_shape.text = slide_heading
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(14)
            run.font.bold = True

    # Tables 1–4 and 6 (excluding CY Returns for now)
    table_data = [key_details, portfolio_aggregates, concentration, top_sectors, cy_returns, top_holdings]
    for i, df in enumerate(table_data, start=1):
        if i == 5:
            continue  # skip CY Returns here
        shape = slide.shapes[i]
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            continue
        table = shape.table
        for row_idx in range(len(df)):
            for col_idx in range(len(df.columns)):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(df.iloc[row_idx, col_idx])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(9)
                        run.font.bold = False

    # CY Returns Table (shape 5)
    cy_table = slide.shapes[5].table

    # Insert header row after the title
    cy_table.cell(1, 0).text = ''
    cy_table.cell(1, 1).text = 'Fund'
    cy_table.cell(1, 2).text = 'Benchmark'
    cy_table.cell(1, 3).text = 'Alpha'

    for col in range(4):
        cell = cy_table.cell(1, col)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(9)
                run.font.bold = True

    # Fill CY Returns content from row 2 onwards
    for row_idx in range(len(cy_returns)):
        for col_idx in range(4):
            cell = cy_table.cell(row_idx + 2, col_idx)
            cell.text = str(cy_returns.iloc[row_idx, col_idx])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(9)
                    run.font.bold = (col_idx == 0)

    """# onwards"""

    list_excel = ["Slide 14 - Parag Flexi", "Slide 15-WOC Flexi", "Slide 16 - Nippon Large", "Slide 17 -SBI Sensex ", "Slide 18 - Nippon Nifty Bees"]
    list_slide = [13, 14, 15, 16, 17]

    for i in range(0, 5):
        df = pd.read_excel(input_path , sheet_name=list_excel[i], header=None)
        slide = prs.slides[list_slide[i]]

        # 1. Key Details
        key_details_heading = df.iloc[3, 1:3].tolist()
        key_details = df.iloc[4:10, 1:3].dropna().reset_index(drop=True)
        key_details.columns = key_details_heading

        # Format Inception date as Aug-22
        if 'Inception' in key_details[key_details.columns[0]].values:
            idx = key_details[key_details[key_details.columns[0]] == 'Inception'].index[0]
            val = key_details.loc[idx, key_details.columns[1]]
            try:
                formatted = pd.to_datetime(str(val), errors='coerce').strftime('%b-%y')
                key_details.loc[idx, key_details.columns[1]] = formatted
            except:
                pass

        # 2. Portfolio Aggregates
        portfolio_aggregates_heading = df.iloc[3, 5:7].tolist()
        portfolio_aggregates = df.iloc[4:11, 5:7].dropna().reset_index(drop=True)
        portfolio_aggregates.columns = portfolio_aggregates_heading
        portfolio_aggregates.iloc[3:, 1] = pd.to_numeric(portfolio_aggregates.iloc[3:, 1], errors='coerce').apply(lambda x: f"{round(x * 100)}%" if pd.notna(x) else "-")

        # 3. Portfolio Concentration
        concentration_heading = df.iloc[3, [9, 11]].tolist()
        concentration_heading = ['Metric', 'Value'] if any(pd.isna(concentration_heading)) else concentration_heading
        concentration = df.iloc[4:9, [9, 11]].dropna().reset_index(drop=True)
        concentration.columns = concentration_heading
        concentration_vals = pd.to_numeric(concentration.iloc[:, 1], errors='coerce')
        concentration.iloc[:, 1] = [
            f"{round(val * 100, 1)}%" if pd.notna(val) and val < 10 else
            f"{round(val, 1)}x" if pd.notna(val) else "-"
            for val in concentration_vals
        ]

        # 4. Top Sectors
        top_sectors_heading = df.iloc[3, 13:15].tolist()
        top_sectors_heading = ['Sector', 'Weight'] if any(pd.isna(top_sectors_heading)) else top_sectors_heading
        top_sectors = df.iloc[4:10, 13:15].dropna().reset_index(drop=True)
        top_sectors.columns = top_sectors_heading
        top_sectors.iloc[:, 1] = pd.to_numeric(top_sectors.iloc[:, 1], errors='coerce').apply(lambda x: f"{round(x * 100, 1)}%" if pd.notna(x) else "-")

        # 5. CY Returns
        cy_returns_heading = df.iloc[17, 1:5].tolist()
        cy_returns = df.iloc[18:27, 1:5].dropna().reset_index(drop=True)
        cy_returns.columns = cy_returns_heading
        for col in cy_returns.columns[1:]:
            cy_returns[col] = pd.to_numeric(cy_returns[col], errors='coerce').apply(lambda x: f"{round(x * 100, 1)}%" if pd.notna(x) else "-")

        # 6. Top Holdings
        top_holdings_heading = df.iloc[16, [7, 8]].tolist()
        top_holdings_heading = ['Top Holdings', '% Allocation'] if any(pd.isna(top_holdings_heading)) else top_holdings_heading
        top_holdings = df.iloc[17:25, [7, 8]].dropna().reset_index(drop=True)
        top_holdings.columns = top_holdings_heading
        top_holdings.iloc[:, 1] = pd.to_numeric(top_holdings.iloc[:, 1], errors='coerce').apply(lambda x: round(x, 2) if pd.notna(x) else "-")

        # Slide Heading
        slide_heading = df.iloc[1, 1]
        title_shape = slide.shapes[0]
        title_shape.text = slide_heading
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
                run.font.bold = True

        # Tables 1–4 & 6
        table_data = [key_details, portfolio_aggregates, concentration, top_sectors, cy_returns, top_holdings]
        for j, data in enumerate(table_data, start=1):
            if j == 5:
                continue
            shape = slide.shapes[j]
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue
            table = shape.table
            for row_idx in range(len(data)):
                for col_idx in range(len(data.columns)):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(data.iloc[row_idx, col_idx])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.name = 'Calibri'
                            run.font.size = Pt(9)
                            run.font.bold = False

        # CY Returns Table (shape 5)
        cy_table = slide.shapes[5].table
        cy_table.cell(1, 0).text = ''
        cy_table.cell(1, 1).text = 'Fund'
        cy_table.cell(1, 2).text = 'Benchmark'
        cy_table.cell(1, 3).text = 'Alpha'

        for col in range(4):
            cell = cy_table.cell(1, col)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(9)
                    run.font.bold = True

        for row_idx in range(len(cy_returns)):
            for col_idx in range(4):
                cell = cy_table.cell(row_idx + 2, col_idx)
                cell.text = str(cy_returns.iloc[row_idx, col_idx])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(9)
                        run.font.bold = (col_idx == 0)

    """#download"""

    output_path = os.path.join(os.path.dirname(__file__), "ppt.pptx")
    prs.save(output_path)
    return output_path
