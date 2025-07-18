from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.table import _Cell, _Row
import copy
import pandas as pd
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

owner_no = input("Enter Owner Number: ")
excel_path = "/content/Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Current Portfolio", header=None)
ppt_path = "/content/Quarterly Review Revised Template_new (1).pptx"
prs = Presentation(ppt_path)

def format_inr_crore(value):
    try:
        value = float(value) / 1e7  # Convert to crores
        return "{:,.2f}".format(value)
    except Exception:
        return "0.00"

"""## Slide 8 & 9"""

## Slide 8
slide = prs.slides[7] #slide 8(index 7)
owner_rows = df[df[0].astype(str).str.contains(f"Owner : {owner_no}")]
if owner_rows.empty:
    raise ValueError("Owner not found.")

owner_row_idx = owner_rows.index[0]
table_start = owner_row_idx + 2

# Dynamically find end of table
table_rows = []
for idx in range(table_start, len(df)):
    row = df.iloc[idx]
    if pd.isna(row[1]) or str(row[0]).startswith("Owner :"):
        break
    table_rows.append(idx)
if not table_rows:
    raise ValueError("No asset allocation data found for this owner.")

table_df = df.iloc[table_rows, 1:4]
table_df.columns = ["Description", "Market Value", "% Assets"]
table_df["% Assets"] = table_df["% Assets"] * 100

for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        num_table_rows = table.rows.__len__()  # total rows including header
        data_rows = min(len(table_df), num_table_rows - 1)  # exclude header row
        for i in range(data_rows):
            row = table_df.iloc[i]
            table.cell(i+1, 0).text = str(row["Description"])
            table.cell(i+1, 1).text = f"{int(row['Market Value']):,}"
            table.cell(i+1, 2).text = f"{row['% Assets']:.2f}%"
            for j in range(3):
                cell = table.cell(i+1, j)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Aptos"
                        run.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER
        break

for shape in slide.shapes:
    if shape.has_chart:
        chart = shape.chart
        chart_data = CategoryChartData()
        chart_data.categories = table_df["Description"]
        chart_data.add_series('Asset Allocation', table_df["% Assets"])
        chart.replace_data(chart_data)
        break

## Slide 9

# Extract relevant values
equity = float(table_df.loc[table_df["Description"].str.strip().str.lower() == "equity", "% Assets"].values[0])
hybrid = float(table_df.loc[table_df["Description"].str.strip().str.lower() == "hybrid", "% Assets"].values[0])
alternative = float(table_df.loc[table_df["Description"].str.strip().str.lower() == "alternative", "% Assets"].values[0])
debt = float(table_df.loc[table_df["Description"].str.strip().str.lower() == "debt", "% Assets"].values[0])

# Apply your formula
net_equity_exposure = equity + 0.5 * hybrid + alternative
net_debt_exposure = 100 - net_equity_exposure

# Strategic allocations (example values)
strategic_equity = 70.0
strategic_debt = 30.0

# Deviations
deviation_equity = net_equity_exposure - strategic_equity
deviation_debt = net_debt_exposure - strategic_debt

slide = prs.slides[8]  # Adjust slide index as needed

for shape in slide.shapes:
    if shape.has_chart:
        chart = shape.chart
        chart_data = CategoryChartData()
        chart_data.categories = ["Net Equity Exposure", "Net Debt Exposure"]
        chart_data.add_series("Current Allocation", [net_equity_exposure, net_debt_exposure])
        chart_data.add_series("Strategic Allocation", [strategic_equity, strategic_debt])
        chart_data.add_series("Deviation", [deviation_equity, deviation_debt])
        chart.replace_data(chart_data)
        break

prs.save("output_filled.pptx")

"""### Slide 12"""

# List of asset vehicles in the order they appear in your PPT
asset_vehicles = [
    "Equity MF", "Equity PMS", "Equity AIF", "Debt MF", "Structured Debt",
    "Hybrid MF", "Alternate Funds", "Gold", "Liquid MF", "Liquid PMS"
]

# Normalize asset vehicle names for matching
asset_vehicles_norm = [av.replace(' ', '').lower() for av in asset_vehicles]

# Load Excel data
excel_path = "/content/Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance 1 Quarter", header=5)

# Normalize column names: remove all spaces and lowercase
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()

for col in ['beginmv', 'endmv', 'inflow', 'outflow', 'profit']:
    if col in df.columns:
        df[col] = df[col].astype(str).str.replace(',', '').replace('nan', '0')
        df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

# Normalize first column for matching
first_col_norm = df.iloc[:, 0].astype(str).str.replace(' ', '', regex=False).str.lower()

# Find the table in the loaded PowerPoint slide (slide 12, 0-based index 11)
slide = prs.slides[11]
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

def format_inr_crore(value):
    try:
        value = float(value) / 1e7  # Convert to crores
        return "{:,.2f}".format(value)
    except Exception:
        return "0.00"

# Prepare to accumulate totals
totals = [0.0] * 7  # For begin_mv, end_mv, growth, cash_flow, net_gain, abs_gain, bench_gain

for i, asset_vehicle in enumerate(asset_vehicles, start=2):
    asset_vehicle_norm = asset_vehicle.replace(' ', '').lower()
    mask = first_col_norm == asset_vehicle_norm
    if mask.any():
        row = df[mask].iloc[0]
        begin_mv = row.get("beginmv", 0)
        end_mv = row.get("endmv", 0)
        growth = end_mv - begin_mv
        inflow = row.get("inflow", 0)
        outflow = row.get("outflow", 0)
        cash_flow = - inflow + outflow
        net_gain = row.get("profit", 0)
        abs_gain = ((end_mv / begin_mv) - 1) * 100 if begin_mv else 0
        # Read Benchmark Gain % from column O (index 14)
        bench_gain = row.iloc[14] if len(row) > 14 else 0
        benchmark = row.get("benchmark", "")
    else:
        begin_mv = end_mv = growth = inflow = outflow = cash_flow = net_gain = abs_gain = bench_gain = 0
        benchmark = ""
    # Prepare values for PPT columns
    values = [
        begin_mv, end_mv, growth, cash_flow, net_gain, abs_gain, bench_gain, benchmark
    ]
    # Add to totals (except benchmark)
    for idx in range(7):
        try:
            totals[idx] += float(values[idx])
        except Exception:
            pass
    # Indices of columns that are monetary values (to be formatted in Cr)
    cr_indices = [0, 1, 2, 3, 4]  # Mar-YY, Jun-YY, Growth, Cash Flow, Net Gain
    for j, value in enumerate(values, start=2):
        cell = table.cell(i, j)
        if j-2 in cr_indices:
            cell.text = format_inr_crore(value)
        elif j-2 == 6:  # Benchmark Gain % (from column O)
            try:
                cell.text = f"{float(value):.2f}"
            except (ValueError, TypeError):
                cell.text = str(value)
        elif j-2 == 5:  # Absolute Gain %
            try:
                cell.text = f"{float(value):.2f}"
            except (ValueError, TypeError):
                cell.text = str(value)
        else:
            cell.text = str(value)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER

# Fill the "Total" row (assumed to be row 12, i=12)
total_row_idx = len(asset_vehicles) + 2
for j, value in enumerate(totals + [""]):  # Add empty string for benchmark column
    cell = table.cell(total_row_idx, j+2)
    if j in [0, 1, 2, 3, 4]:
        cell.text = format_inr_crore(value)
    elif j in [5, 6]:
        try:
            cell.text = f"{float(value):.2f}"
        except (ValueError, TypeError):
            cell.text = str(value)
    else:
        cell.text = ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER

prs.save("output_filled.pptx")

"""### Slide 13"""

# List of asset vehicles in the order they appear in your PPT
asset_vehicles = [
    "Equity MF", "Equity PMS", "Equity AIF", "Debt MF", "Structured Debt",
    "Hybrid MF", "Alternate Funds", "Gold", "Liquid MF", "Liquid PMS"
]

# Normalize asset vehicle names for matching
asset_vehicles_norm = [av.replace(' ', '').lower() for av in asset_vehicles]

# Load Excel data
excel_path = "/content/Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance 1 year", header=5)

# Normalize column names: remove all spaces and lowercase
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()

for col in ['beginmv', 'endmv', 'inflow', 'outflow', 'profit']:
    if col in df.columns:
        df[col] = df[col].astype(str).str.replace(',', '').replace('nan', '0')
        df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

# Normalize first column for matching
first_col_norm = df.iloc[:, 0].astype(str).str.replace(' ', '', regex=False).str.lower()

# Load PowerPoint
slide = prs.slides[12]  # Slide 13 (0-based index)

# Find the table
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

def format_inr_crore(value):
    try:
        value = float(value) / 1e7  # Convert to crores
        return "{:,.2f}".format(value)
    except Exception:
        return "0.00"

# Prepare to accumulate totals
totals = [0.0] * 7  # For begin_mv, end_mv, growth, cash_flow, net_gain, abs_gain, bench_gain

for i, asset_vehicle in enumerate(asset_vehicles, start=2):
    asset_vehicle_norm = asset_vehicle.replace(' ', '').lower()
    mask = first_col_norm == asset_vehicle_norm
    if mask.any():
        row = df[mask].iloc[0]
        begin_mv = row.get("beginmv", 0)
        end_mv = row.get("endmv", 0)
        growth = end_mv - begin_mv
        inflow = row.get("inflow", 0)
        outflow = row.get("outflow", 0)
        cash_flow = - inflow + outflow
        net_gain = row.get("profit", 0)
        abs_gain = ((end_mv / begin_mv) - 1) * 100 if begin_mv else 0
        # Read Benchmark Gain % from column U (index 20)
        bench_gain = row.iloc[20] if len(row) > 20 else 0
        benchmark = row.get("benchmark", "")
    else:
        begin_mv = end_mv = growth = inflow = outflow = cash_flow = net_gain = abs_gain = bench_gain = 0
        benchmark = ""
    # Prepare values for PPT columns
    values = [
        begin_mv, end_mv, growth, cash_flow, net_gain, abs_gain, bench_gain, benchmark
    ]
    # Add to totals (except benchmark)
    for idx in range(7):
        try:
            totals[idx] += float(values[idx])
        except Exception:
            pass
    # Indices of columns that are monetary values (to be formatted in Cr)
    cr_indices = [0, 1, 2, 3, 4]  # Mar-YY, Jun-YY, Growth, Cash Flow, Net Gain
    for j, value in enumerate(values, start=2):
        cell = table.cell(i, j)
        if j-2 in cr_indices:
            cell.text = format_inr_crore(value)
        elif j-2 == 6:  # Benchmark Gain % (from column U)
            try:
                cell.text = f"{float(value):.2f}"
            except (ValueError, TypeError):
                cell.text = str(value)
        elif j-2 == 5:  # Absolute Gain %
            try:
                cell.text = f"{float(value):.2f}"
            except (ValueError, TypeError):
                cell.text = str(value)
        else:
            cell.text = str(value)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER

# Fill the "Total" row (assumed to be row 12, i=12)
total_row_idx = len(asset_vehicles) + 2
for j, value in enumerate(totals + [""]):  # Add empty string for benchmark column
    cell = table.cell(total_row_idx, j+2)
    if j in [0, 1, 2, 3, 4]:
        cell.text = format_inr_crore(value)
    elif j in [5, 6]:
        try:
            cell.text = f"{float(value):.2f}"
        except (ValueError, TypeError):
            cell.text = str(value)
    else:
        cell.text = ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER

prs.save("output_filled.pptx")

"""## Slide 14"""

excel_path = "Automation Reports.xlsx"
owner_no = "11870"  # Set your owner number here

# --- 1. Find Owner Row and Extract Table from Current Portfolio ---
df_portfolio = pd.read_excel(excel_path, sheet_name="Current Portfolio", header=None)
owner_rows = df_portfolio[df_portfolio[0].astype(str).str.contains(f"Owner : {owner_no}")]
if owner_rows.empty:
    raise ValueError("Owner not found.")

owner_row_idx = owner_rows.index[0]
table_start = owner_row_idx + 2

# Dynamically find end of table
table_rows = []
for idx in range(table_start, len(df_portfolio)):
    row = df_portfolio.iloc[idx]
    if pd.isna(row[1]) or str(row[0]).startswith("Owner :"):
        break
    table_rows.append(idx)
if not table_rows:
    raise ValueError("No asset allocation data found for this owner.")

table_df = df_portfolio.iloc[table_rows, 1:4]
table_df.columns = ["Description", "Market Value", "% Assets"]
table_df["% Assets"] = pd.to_numeric(table_df["% Assets"], errors="coerce").fillna(0) * 100

# --- 2. Load Performance SI, 1 Quarter, 1 Year for this owner ---
def load_owner_df(sheet, owner_no):
    df = pd.read_excel(excel_path, sheet_name=sheet, header=5)
    df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()
    if 'owner' in df.columns:
        df['owner'] = df['owner'].astype(str)
        return df[df['owner'].str.contains(owner_no, na=False)]
    return df

df_si = load_owner_df("Performance SI", owner_no)
df_qtr = load_owner_df("Performance 1 Quarter", owner_no)
df_1yr = load_owner_df("Performance 1 year", owner_no)

# --- 3. Asset Vehicles as per PPT order ---
asset_vehicles = [
    "Equity MF", "Equity PMS", "Equity AIF", "Debt MF", "Structured Debt",
    "Hybrid MF", "Alternate Funds", "Gold", "Liquid MF", "Liquid PMS", "Total"
]
asset_vehicles_norm = [av.replace(' ', '').lower() for av in asset_vehicles]

# --- 4. Helper Functions ---
def format_inr_cr(val):
    try:
        return f"{float(val)/1e7:,.2f}"
    except Exception:
        return ""

def format_pct(val):
    try:
        return f"{float(val):.2f}%"
    except Exception:
        return ""

def format_xirr(val):
    try:
        v = float(str(val).replace('%','').replace(',','').strip())
        if abs(v) < 1.5:
            v *= 100
        return f"{v:.2f}%"
    except Exception:
        return ""

def get_benchmark_from_si(df_si, sec_norm):
    match = df_si[df_si['security'].str.replace(' ', '').str.lower() == sec_norm]
    if not match.empty and "benchmark" in match.columns:
        return match.iloc[0]['benchmark']
    return ""

# --- 5. Prepare Data for Each Asset Vehicle ---
rows = []
total_invested = 0
total_market_value = 0

for sec_type in asset_vehicles:
    sec_norm = sec_type.replace(' ', '').lower()
    # Amount Invested (Total Cost) from Current Portfolio
    mask_portfolio = table_df["Description"].astype(str).str.replace(' ', '').str.lower() == sec_norm
    amount_invested = table_df[mask_portfolio]["Market Value"].sum() if not table_df[mask_portfolio].empty else 0
    amount_invested = float(amount_invested) if amount_invested != "" else 0.0

    # Market Value (End MV) from Performance SI
    mask_si = df_si['security'].str.replace(' ', '').str.lower() == sec_norm
    try:
      market_value_clean = str(market_value).replace(',', '').strip()
      market_value = float(market_value_clean) if market_value_clean != "" else 0.0
    except Exception:
      market_value = 0.0


    total_invested += amount_invested
    total_market_value += market_value

    # Portfolio (Qx'YY) and Benchmark (Qx'YY) from Performance 1 Quarter (row 7: index 6)
    mask_qtr = df_qtr['security'].str.replace(' ', '').str.lower() == sec_norm
    port_qxyy = df_qtr[mask_qtr].iloc[0, 11] if not df_qtr[mask_qtr].empty and len(df_qtr.columns) > 11 else ""
    bench_qxyy = df_qtr[mask_qtr].iloc[0, 17] if not df_qtr[mask_qtr].empty and len(df_qtr.columns) > 17 else ""

    # Portfolio XIRR (1 Yr.) and Benchmark XIRR (1 Yr.) from Performance 1 year (row 7: index 6)
    mask_1yr = df_1yr['security'].str.replace(' ', '').str.lower() == sec_norm
    port_xirr_1yr = format_xirr(df_1yr[mask_1yr].iloc[0, 17]) if not df_1yr[mask_1yr].empty and len(df_1yr.columns) > 17 else ""
    bench_xirr_1yr = format_xirr(df_1yr[mask_1yr].iloc[0, 20]) if not df_1yr[mask_1yr].empty and len(df_1yr.columns) > 20 else ""

    # Portfolio XIRR (SI) and Benchmark XIRR (SI) from Performance SI (row 7: index 6)
    port_xirr_si = format_xirr(df_si[mask_si].iloc[0, 17]) if not df_si[mask_si].empty and len(df_si.columns) > 17 else ""
    bench_xirr_si = format_xirr(df_si[mask_si].iloc[0, 20]) if not df_si[mask_si].empty and len(df_si.columns) > 20 else ""

    # Benchmark from Performance SI
    benchmark = get_benchmark_from_si(df_si, sec_norm)

    # Allocation % (will fill after total is known)
    rows.append([
        sec_type,
        amount_invested,
        0,  # Placeholder for allocation %
        market_value,
        port_qxyy,
        bench_qxyy,
        port_xirr_1yr,
        bench_xirr_1yr,
        port_xirr_si,
        bench_xirr_si,
        benchmark
    ])

# --- 6. Calculate Allocation % ---
for row in rows:
    invested = row[1]
    row[2] = format_pct((invested / total_invested * 100) if total_invested else 0)

# --- 7. Fill the PowerPoint Table ---
slide = prs.slides[13]  # Slide 14 (0-based index)

# Find the table
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on slide 8!")

# Fill the table
for i, row in enumerate(rows):
    table.cell(i+1, 0).text = str(row[0])  # Security Type
    table.cell(i+1, 1).text = format_inr_cr(row[1])  # Amount Invested
    table.cell(i+1, 2).text = row[2]  # Allocation %
    table.cell(i+1, 3).text = format_inr_cr(row[3])  # Market Value
    table.cell(i+1, 4).text = str(row[4])  # Portfolio (Qx'YY)
    table.cell(i+1, 5).text = str(row[5])  # Benchmark (Qx'YY)
    table.cell(i+1, 6).text = str(row[6])  # Portfolio XIRR (1 Yr.)
    table.cell(i+1, 7).text = str(row[7])  # Benchmark XIRR (1 Yr.)
    table.cell(i+1, 8).text = str(row[8])  # Portfolio XIRR (SI)
    table.cell(i+1, 9).text = str(row[9])  # Benchmark XIRR (SI)
    table.cell(i+1, 10).text = str(row[10])  # Benchmark
    for j in range(11):
        cell = table.cell(i+1, j)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER

prs.save("output_filled.pptx")
print("Slide 14 table filled and saved as output_filled.pptx")

"""## Slide 16"""

# --- 1. Load Excel data ---
excel_path = "Portfolio-performance for wt avg days.xlsx"
df = pd.read_excel(excel_path)

# Normalize column names for matching
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()

# Filter for ASTCLS == 'equity'
df = df[df['astcls'].astype(str).str.strip().str.lower() == 'equity']

# Convert Market Value to numeric (handle commas, NaN, blanks)
if 'marketvalue' not in df.columns:
    raise Exception("Column 'marketvalue' not found in the sheet.")
df['marketvalue'] = pd.to_numeric(
    df['marketvalue'].astype(str).str.replace(',', '').replace('nan', ''), errors='coerce'
).fillna(0)

# Group by Asset Class (Sector Name or similar column)
asset_class_col = 'sectorname' if 'sectorname' in df.columns else 'sector name'
df[asset_class_col] = df[asset_class_col].astype(str).str.strip()

# Calculate sum of Market Value for each asset class
df_summary = df.groupby(asset_class_col)['marketvalue'].sum().reset_index()

# Format Market Value as INR Cr (optional)
df_summary['marketvalue_cr'] = df_summary['marketvalue'] / 1e7
df_summary['marketvalue_cr'] = df_summary['marketvalue_cr'].apply(lambda x: f"{x:,.2f}")

# --- 2. Open PowerPoint and find the table and chart ---
slide = prs.slides[15]  # Replace with the correct slide index (0-based)

# Find the first table on the slide
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on the target slide!")

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

# --- 3. Fill the table with asset class and Market Value sum ---
for i, (_, row) in enumerate(df_summary.iterrows(), start=1):
    if i >= len(table.rows):
        break
    set_cell(table.cell(i, 0), row[asset_class_col], bold=False)
    set_cell(table.cell(i, 1), row['marketvalue_cr'], bold=False)

# Optionally fill a "Total" row
if len(df_summary) + 1 < len(table.rows):
    set_cell(table.cell(len(df_summary)+1, 0), "Total", bold=True)
    set_cell(table.cell(len(df_summary)+1, 1), f"{df_summary['marketvalue'].sum()/1e7:,.2f}", bold=True)

# --- 4. Update the chart for asset class allocation ---
chart_data = CategoryChartData()
chart_data.categories = list(df_summary[asset_class_col])
chart_data.add_series('Market Value', [float(x.replace(',', '')) for x in df_summary['marketvalue_cr']])

for shape in slide.shapes:
    if hasattr(shape, "has_chart") and shape.has_chart:
        chart = shape.chart
        chart.replace_data(chart_data)
        break

prs.save("output_filled.pptx")
print("PowerPoint table and chart filled and saved as output_filled.pptx")

"""## Slide 17

"""

# 1. Load Excel data
excel_path = "Portfolio-performance for wt avg days.xlsx"
df = pd.read_excel(excel_path)
df = df[df['ASTCLS'].str.strip().str.lower() == 'equity'].copy()

# 2. Compute total equity amount invested (TOTAL COST) for allocation %
df['TOTAL COST'] = pd.to_numeric(df['TOTAL COST'], errors='coerce').fillna(0)
total_equity_cost = df['TOTAL COST'].sum()

def calc_allocation_pct(tc):
    try:
        return f"{(float(tc)/total_equity_cost)*100:.2f}%"
    except Exception:
        return ""

# 3. Formatting functions
def format_wa_days(value):
    try:
        return str(int(round(float(value))))
    except Exception:
        return ""

def format_xirr(value):
    try:
        value = str(value).replace('%', '').strip()
        if value == "":
            return ""
        value = float(value)
        if abs(value) > 1.5:
            return f"{value:.2f}%"
        else:
            return f"{value*100:.2f}%"
    except Exception:
        return ""

def format_market_value(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

def format_total_cost(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

# 4. Load PBS Quarter for Alternatives Breakup sheet
pbs_path = "/content/PBS Quarter for Alternatives Breakup.csv"
pbs_df = pd.read_csv(pbs_path, header=2)
pbs_df['Security'] = pbs_df['Security'].astype(str).str.strip().str.lower()

pbs_lookup = {}
for idx, row in pbs_df.iterrows():
    fund = row['Security']
    portfolio_val = row.iloc[6] if len(row) > 6 else ""
    benchmark_val = row.iloc[7] if len(row) > 7 else ""
    benchmark_name = row.get('Benchmark', "")
    pbs_lookup[fund] = {
        'scheme_qxyy': portfolio_val,
        'benchmark_qxyy': benchmark_val,
        'benchmark': benchmark_name
    }

# 5. Open PowerPoint and find table
slide = prs.slides[16]  # Slide 17 (0-indexed)

table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on slide 17!")

def add_row(table):
    new_row = copy.deepcopy(table._tbl.tr_lst[-1])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''
    table._tbl.append(new_row)
    return _Row(new_row, table)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

columns_map = [
    'FUND NAME',           # 0 Scheme Name
    'ALLOCATION_PCT',      # 1 Allocation % (calculated)
    'TOTAL COST',          # 2 Amount Invested
    'WA DAYS',             # 3 WA Days
    'MARKET VALUE',        # 4 Market Value
    'scheme_qxyy',         # 5 Scheme (Qx’YY) from PBS
    'benchmark_qxyy',      # 6 Benchmark (Qx’YY) from PBS
    'XIRR SCHEME (1Y)',    # 7 Scheme XIRR (1 Yr.)
    'XIRR BENCHMARK (1Y)', # 8 Benchmark XIRR (1 Yr.)
    'XIRR SCHEME (SI)',    # 9 Scheme XIRR (SI)
    'XIRR BENCHMARK (SI)', # 10 Benchmark XIRR (SI)
    'benchmark'            # 11 Benchmark (from PBS)
]

# 6. Group by Sector Name and fill the table
row_idx = 1  # Start after header
sector_order = df['SECTOR NAME'].dropna().unique()

for sector in sector_order:
    sector_funds = df[df['SECTOR NAME'] == sector]
    # Add sector heading row (bold)
    if row_idx >= len(table.rows):
        add_row(table)
    set_cell(table.cell(row_idx, 0), sector, bold=True)
    for col in range(1, len(table.columns)):
        set_cell(table.cell(row_idx, col), "", bold=True)
    row_idx += 1

    # Add fund rows for this sector
    for _, fund_row in sector_funds.iterrows():
        if row_idx >= len(table.rows):
            add_row(table)
        fund_name = str(fund_row['FUND NAME']).strip().lower()
        for col_idx, excel_col in enumerate(columns_map):
            value = ""
            if excel_col:
                if excel_col == 'ALLOCATION_PCT':
                    value = calc_allocation_pct(fund_row['TOTAL COST'])
                elif excel_col == 'TOTAL COST':
                    value = format_total_cost(fund_row['TOTAL COST'])
                elif excel_col == 'WA DAYS':
                    value = format_wa_days(fund_row['WA DAYS'])
                elif excel_col == 'MARKET VALUE':
                    value = format_market_value(fund_row['MARKET VALUE'])
                elif excel_col in ['XIRR SCHEME (1Y)', 'XIRR BENCHMARK (1Y)', 'XIRR SCHEME (SI)', 'XIRR BENCHMARK (SI)']:
                    value = format_xirr(fund_row[excel_col])
                elif excel_col in ['scheme_qxyy', 'benchmark_qxyy', 'benchmark']:
                    value = pbs_lookup.get(fund_name, {}).get(excel_col, "")
                else:
                    value = fund_row.get(excel_col, "")
                set_cell(table.cell(row_idx, col_idx), value, bold=False)
            else:
                set_cell(table.cell(row_idx, col_idx), "", bold=False)
        row_idx += 1

prs.save("output_filled.pptx")
print("PowerPoint table filled and saved as output_filled.pptx")

"""## Slide 20"""

# 1. Load Excel data
excel_path = "Portfolio-performance for wt avg days.xlsx"
df = pd.read_excel(excel_path)
df = df[df['ASTCLS'].str.strip().str.lower() == 'debt'].copy()

# 2. Compute total equity amount invested (TOTAL COST) for allocation %
df['TOTAL COST'] = pd.to_numeric(df['TOTAL COST'], errors='coerce').fillna(0)
total_equity_cost = df['TOTAL COST'].sum()

def calc_allocation_pct(tc):
    try:
        return f"{(float(tc)/total_equity_cost)*100:.2f}%"
    except Exception:
        return ""

# 3. Formatting functions
def format_wa_days(value):
    try:
        return str(int(round(float(value))))
    except Exception:
        return ""

def format_xirr(value):
    try:
        value = str(value).replace('%', '').strip()
        if value == "":
            return ""
        value = float(value)
        if abs(value) > 1.5:
            return f"{value:.2f}%"
        else:
            return f"{value*100:.2f}%"
    except Exception:
        return ""

def format_market_value(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

def format_total_cost(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

# 4. Load PBS Quarter for Alternatives Breakup sheet
pbs_path = "/content/PBS Quarter for Alternatives Breakup.csv"
pbs_df = pd.read_csv(pbs_path, header=2)
pbs_df['Security'] = pbs_df['Security'].astype(str).str.strip().str.lower()

pbs_lookup = {}
for idx, row in pbs_df.iterrows():
    fund = row['Security']
    portfolio_val = row.iloc[6] if len(row) > 6 else ""
    benchmark_val = row.iloc[7] if len(row) > 7 else ""
    benchmark_name = row.get('Benchmark', "")
    pbs_lookup[fund] = {
        'scheme_qxyy': portfolio_val,
        'benchmark_qxyy': benchmark_val,
        'benchmark': benchmark_name
    }

# 5. Open PowerPoint and find table
slide = prs.slides[19]  # Slide 20 (0-indexed)

table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on slide 20!")

def add_row(table):
    new_row = copy.deepcopy(table._tbl.tr_lst[-1])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''
    table._tbl.append(new_row)
    return _Row(new_row, table)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

columns_map = [
    'FUND NAME',           # 0 Scheme Name
    'ALLOCATION_PCT',      # 1 Allocation % (calculated)
    'TOTAL COST',          # 2 Amount Invested
    'WA DAYS',             # 3 WA Days
    'MARKET VALUE',        # 4 Market Value
    'scheme_qxyy',         # 5 Scheme (Qx’YY) from PBS
    'benchmark_qxyy',      # 6 Benchmark (Qx’YY) from PBS
    'XIRR SCHEME (1Y)',    # 7 Scheme XIRR (1 Yr.)
    'XIRR BENCHMARK (1Y)', # 8 Benchmark XIRR (1 Yr.)
    'XIRR SCHEME (SI)',    # 9 Scheme XIRR (SI)
    'XIRR BENCHMARK (SI)', # 10 Benchmark XIRR (SI)
    'benchmark'            # 11 Benchmark (from PBS)
]

# 6. Group by Sector Name and fill the table
row_idx = 1  # Start after header
sector_order = df['SECTOR NAME'].dropna().unique()

for sector in sector_order:
    sector_funds = df[df['SECTOR NAME'] == sector]
    # Add sector heading row (bold)
    if row_idx >= len(table.rows):
        add_row(table)
    set_cell(table.cell(row_idx, 0), sector, bold=True)
    for col in range(1, len(table.columns)):
        set_cell(table.cell(row_idx, col), "", bold=True)
    row_idx += 1

    # Add fund rows for this sector
    for _, fund_row in sector_funds.iterrows():
        if row_idx >= len(table.rows):
            add_row(table)
        fund_name = str(fund_row['FUND NAME']).strip().lower()
        for col_idx, excel_col in enumerate(columns_map):
            value = ""
            if excel_col:
                if excel_col == 'ALLOCATION_PCT':
                    value = calc_allocation_pct(fund_row['TOTAL COST'])
                elif excel_col == 'TOTAL COST':
                    value = format_total_cost(fund_row['TOTAL COST'])
                elif excel_col == 'WA DAYS':
                    value = format_wa_days(fund_row['WA DAYS'])
                elif excel_col == 'MARKET VALUE':
                    value = format_market_value(fund_row['MARKET VALUE'])
                elif excel_col in ['XIRR SCHEME (1Y)', 'XIRR BENCHMARK (1Y)', 'XIRR SCHEME (SI)', 'XIRR BENCHMARK (SI)']:
                    value = format_xirr(fund_row[excel_col])
                elif excel_col in ['scheme_qxyy', 'benchmark_qxyy', 'benchmark']:
                    value = pbs_lookup.get(fund_name, {}).get(excel_col, "")
                else:
                    value = fund_row.get(excel_col, "")
                set_cell(table.cell(row_idx, col_idx), value, bold=False)
            else:
                set_cell(table.cell(row_idx, col_idx), "", bold=False)
        row_idx += 1

prs.save("output_filled.pptx")
print("PowerPoint table filled and saved as output_filled.pptx")

"""## Slide 24"""

# 1. Load Excel data
excel_path = "Portfolio-performance for wt avg days.xlsx"
df = pd.read_excel(excel_path)
df = df[df['ASTCLS'].str.strip().str.lower() == 'hybrid'].copy()

# 2. Compute total equity amount invested (TOTAL COST) for allocation %
df['TOTAL COST'] = pd.to_numeric(df['TOTAL COST'], errors='coerce').fillna(0)
total_equity_cost = df['TOTAL COST'].sum()

def calc_allocation_pct(tc):
    try:
        return f"{(float(tc)/total_equity_cost)*100:.2f}%"
    except Exception:
        return ""

# 3. Formatting functions
def format_wa_days(value):
    try:
        return str(int(round(float(value))))
    except Exception:
        return ""

def format_xirr(value):
    try:
        value = str(value).replace('%', '').strip()
        if value == "":
            return ""
        value = float(value)
        if abs(value) > 1.5:
            return f"{value:.2f}%"
        else:
            return f"{value*100:.2f}%"
    except Exception:
        return ""

def format_market_value(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

def format_total_cost(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

# 4. Load PBS Quarter for Alternatives Breakup sheet
pbs_path = "/content/PBS Quarter for Alternatives Breakup.csv"
pbs_df = pd.read_csv(pbs_path, header=2)
pbs_df['Security'] = pbs_df['Security'].astype(str).str.strip().str.lower()

pbs_lookup = {}
for idx, row in pbs_df.iterrows():
    fund = row['Security']
    portfolio_val = row.iloc[6] if len(row) > 6 else ""
    benchmark_val = row.iloc[7] if len(row) > 7 else ""
    benchmark_name = row.get('Benchmark', "")
    pbs_lookup[fund] = {
        'scheme_qxyy': portfolio_val,
        'benchmark_qxyy': benchmark_val,
        'benchmark': benchmark_name
    }

# 5. Open PowerPoint and find table
slide = prs.slides[23]  # Slide 24 (0-indexed)

table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on slide 14!")

def add_row(table):
    new_row = copy.deepcopy(table._tbl.tr_lst[-1])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''
    table._tbl.append(new_row)
    return _Row(new_row, table)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

columns_map = [
    'FUND NAME',           # 0 Scheme Name
    'ALLOCATION_PCT',      # 1 Allocation % (calculated)
    'TOTAL COST',          # 2 Amount Invested
    'WA DAYS',             # 3 WA Days
    'MARKET VALUE',        # 4 Market Value
    'scheme_qxyy',         # 5 Scheme (Qx’YY) from PBS
    'benchmark_qxyy',      # 6 Benchmark (Qx’YY) from PBS
    'XIRR SCHEME (1Y)',    # 7 Scheme XIRR (1 Yr.)
    'XIRR BENCHMARK (1Y)', # 8 Benchmark XIRR (1 Yr.)
    'XIRR SCHEME (SI)',    # 9 Scheme XIRR (SI)
    'XIRR BENCHMARK (SI)', # 10 Benchmark XIRR (SI)
    'benchmark'            # 11 Benchmark (from PBS)
]

# 6. Group by Sector Name and fill the table
row_idx = 1  # Start after header
sector_order = df['SECTOR NAME'].dropna().unique()

for sector in sector_order:
    sector_funds = df[df['SECTOR NAME'] == sector]
    # Add sector heading row (bold)
    if row_idx >= len(table.rows):
        add_row(table)
    set_cell(table.cell(row_idx, 0), sector, bold=True)
    for col in range(1, len(table.columns)):
        set_cell(table.cell(row_idx, col), "", bold=True)
    row_idx += 1

    # Add fund rows for this sector
    for _, fund_row in sector_funds.iterrows():
        if row_idx >= len(table.rows):
            add_row(table)
        fund_name = str(fund_row['FUND NAME']).strip().lower()
        for col_idx, excel_col in enumerate(columns_map):
            value = ""
            if excel_col:
                if excel_col == 'ALLOCATION_PCT':
                    value = calc_allocation_pct(fund_row['TOTAL COST'])
                elif excel_col == 'TOTAL COST':
                    value = format_total_cost(fund_row['TOTAL COST'])
                elif excel_col == 'WA DAYS':
                    value = format_wa_days(fund_row['WA DAYS'])
                elif excel_col == 'MARKET VALUE':
                    value = format_market_value(fund_row['MARKET VALUE'])
                elif excel_col in ['XIRR SCHEME (1Y)', 'XIRR BENCHMARK (1Y)', 'XIRR SCHEME (SI)', 'XIRR BENCHMARK (SI)']:
                    value = format_xirr(fund_row[excel_col])
                elif excel_col in ['scheme_qxyy', 'benchmark_qxyy', 'benchmark']:
                    value = pbs_lookup.get(fund_name, {}).get(excel_col, "")
                else:
                    value = fund_row.get(excel_col, "")
                set_cell(table.cell(row_idx, col_idx), value, bold=False)
            else:
                set_cell(table.cell(row_idx, col_idx), "", bold=False)
        row_idx += 1

prs.save("output_filled.pptx")
print("PowerPoint table filled and saved as output_filled.pptx")

"""## Slide 27"""

# 1. Load Excel data
excel_path = "Portfolio-performance for wt avg days.xlsx"
df = pd.read_excel(excel_path)
df = df[
    (df['ASTCLS'].str.strip().str.lower() == 'alternative') &
    (df['SECTOR NAME'].str.strip().str.lower() == 'international')
].copy()

# 2. Compute total alternative amount invested (TOTAL COST) for allocation %
df['TOTAL COST'] = pd.to_numeric(df['TOTAL COST'], errors='coerce').fillna(0)
total_alternative_cost = df['TOTAL COST'].sum()

def calc_allocation_pct(tc):
    try:
        return f"{(float(tc)/total_alternative_cost)*100:.2f}%"
    except Exception:
        return ""

def format_wa_days(value):
    try:
        return str(int(round(float(value))))
    except Exception:
        return ""

def format_xirr(value):
    try:
        value = str(value).replace('%', '').strip()
        if value == "":
            return ""
        value = float(value)
        if abs(value) > 1.5:
            return f"{value:.2f}%"
        else:
            return f"{value*100:.2f}%"
    except Exception:
        return ""

def format_market_value(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

def format_total_cost(value):
    try:
        return f"{float(value):,.2f}"
    except Exception:
        return ""

# 3. Load PBS Quarter for Alternatives Breakup sheet
pbs_path = "/content/PBS Quarter for Alternatives Breakup.csv"
pbs_df = pd.read_csv(pbs_path, header=2)
pbs_df['Security'] = pbs_df['Security'].astype(str).str.strip().str.lower()

pbs_lookup = {}
for idx, row in pbs_df.iterrows():
    fund = row['Security']
    portfolio_val = row.iloc[6] if len(row) > 6 else ""
    benchmark_val = row.iloc[7] if len(row) > 7 else ""
    benchmark_name = row.get('Benchmark', "")
    pbs_lookup[fund] = {
        'scheme_qxyy': portfolio_val,
        'benchmark_qxyy': benchmark_val,
        'benchmark': benchmark_name
    }

# 4. Open PowerPoint and find table
slide = prs.slides[26]  # Slide 27 (0-indexed)

table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break
if table is None:
    raise Exception("No table found on slide 27!")

def add_row(table):
    new_row = copy.deepcopy(table._tbl.tr_lst[-1])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''
    table._tbl.append(new_row)
    return _Row(new_row, table)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

columns_map = [
    'FUND NAME',           # 0 Scheme Name
    'ALLOCATION_PCT',      # 1 Allocation % (calculated)
    'TOTAL COST',          # 2 Amount Invested
    'WA DAYS',             # 3 WA Days
    'MARKET VALUE',        # 4 Market Value
    'scheme_qxyy',         # 5 Scheme (Qx’YY) from PBS
    'benchmark_qxyy',      # 6 Benchmark (Qx’YY) from PBS
    'XIRR SCHEME (1Y)',    # 7 Scheme XIRR (1 Yr.)
    'XIRR BENCHMARK (1Y)', # 8 Benchmark XIRR (1 Yr.)
    'XIRR SCHEME (SI)',    # 9 Scheme XIRR (SI)
    'XIRR BENCHMARK (SI)', # 10 Benchmark XIRR (SI)
    'benchmark'            # 11 Benchmark (from PBS)
]

# 5. Fill the table (only International sector funds)
row_idx = 1  # Start after header

sector = "International"
sector_funds = df[df['SECTOR NAME'].str.strip().str.lower() == 'international']

if not sector_funds.empty:
    # Add sector heading row (bold)
    if row_idx >= len(table.rows):
        add_row(table)
    set_cell(table.cell(row_idx, 0), sector, bold=True)
    for col in range(1, len(table.columns)):
        set_cell(table.cell(row_idx, col), "", bold=True)
    row_idx += 1

    # Add fund rows for this sector
    for _, fund_row in sector_funds.iterrows():
        if row_idx >= len(table.rows):
            add_row(table)
        fund_name = str(fund_row['FUND NAME']).strip().lower()
        for col_idx, excel_col in enumerate(columns_map):
            value = ""
            if excel_col:
                if excel_col == 'ALLOCATION_PCT':
                    value = calc_allocation_pct(fund_row['TOTAL COST'])
                elif excel_col == 'TOTAL COST':
                    value = format_total_cost(fund_row['TOTAL COST'])
                elif excel_col == 'WA DAYS':
                    value = format_wa_days(fund_row['WA DAYS'])
                elif excel_col == 'MARKET VALUE':
                    value = format_market_value(fund_row['MARKET VALUE'])
                elif excel_col in ['XIRR SCHEME (1Y)', 'XIRR BENCHMARK (1Y)', 'XIRR SCHEME (SI)', 'XIRR BENCHMARK (SI)']:
                    value = format_xirr(fund_row[excel_col])
                elif excel_col in ['scheme_qxyy', 'benchmark_qxyy', 'benchmark']:
                    value = pbs_lookup.get(fund_name, {}).get(excel_col, "")
                else:
                    value = fund_row.get(excel_col, "")
                set_cell(table.cell(row_idx, col_idx), value, bold=False)
            else:
                set_cell(table.cell(row_idx, col_idx), "", bold=False)
        row_idx += 1

prs.save("output_filled.pptx")
print("PowerPoint table filled and saved as output_filled.pptx")

"""## Slide 18"""

# Load the data
excel_path = "Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance SI", header=5)

# Normalize column names and fund names
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()
df['security'] = df['security'].astype(str).str.strip()

# Filter for relevant rows (between 'Equity' and 'Debt')
start_idx = df[df['security'].str.lower() == 'equity'].index[0]
try:
    end_idx = df[df['security'].str.lower() == 'debt'].index[0]
except IndexError:
    end_idx = len(df)

df_section = df.loc[start_idx+1:end_idx-1].copy()

# Convert End MV to numeric, handling commas and missing values
df_section['endmv'] = pd.to_numeric(df_section['endmv'].astype(str).str.replace(',', ''), errors='coerce').fillna(0)

# Filter funds with End MV > 0
df_filtered = df_section[df_section['endmv'] > 0].copy()

# Calculate Portfolio and Benchmark annualized returns
def safe_float(x):
    try:
        return float(str(x).replace('%', '').strip())
    except:
        return 0.0

df_filtered['portfolio_ann'] = df_filtered.iloc[:, 11].apply(safe_float)  # Column L (index 11)
df_filtered['benchmark_ann'] = df_filtered.iloc[:, 14].apply(safe_float)  # Column O (index 14)

# Calculate outperformance
df_filtered['outperformance'] = df_filtered['portfolio_ann'] - df_filtered['benchmark_ann']

# Sort and get top 5 outperforming and top 1 underperforming
df_outperforming = df_filtered[df_filtered['outperformance'] > 0].sort_values(by='outperformance', ascending=False).head(5)
df_underperforming = df_filtered[df_filtered['outperformance'] <= 0].sort_values(by='outperformance').head(1)

print("Top 5 Outperforming funds:")
print(df_outperforming[['security', 'endmv', 'outperformance']])

print("\nTop 1 Underperforming fund:")
print(df_underperforming[['security', 'endmv', 'outperformance']])

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.table import _Cell, _Row
import copy

# --- 1. Load Excel data ---
excel_path = "Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance SI", header=5)

# Normalize columns and first column for matching
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()
first_col = df.columns[0]
df[first_col] = df[first_col].astype(str).str.strip()

# Find start and stop indices for Equity section in first column
equity_idx = df[df[first_col].str.lower() == 'equity'].index
debt_idx = df[df[first_col].str.lower() == 'debt'].index

if len(equity_idx) == 0:
    raise Exception("No 'Equity' row found in first column.")
if len(debt_idx) == 0:
    stop_idx = len(df)
else:
    stop_idx = debt_idx[0]

start_idx = equity_idx[0]

# Slice the dataframe for Equity schemes only
df_equity = df.loc[start_idx+1:stop_idx-1].copy()

# Skip unwanted scheme names
skip_names = ['equity mf', 'equity pms', 'equity aif', 'total : stock', 'total : equity']
df_equity = df_equity[~df_equity[first_col].str.lower().isin(skip_names)]

# Convert endmv to numeric (handle commas, NaN, blanks)
df_equity['endmv'] = pd.to_numeric(
    df_equity['endmv'].astype(str).str.replace(',', '').replace('nan', ''), errors='coerce'
).fillna(0)

# Only schemes with End MV > 0
df_equity = df_equity[df_equity['endmv'] > 0]

# Calculate total End MV for allocation %
total_mv = df_equity['endmv'].sum()

# Calculate Portfolio and Benchmark Annualized Returns and Outperformance
def safe_float(val):
    try:
        return float(str(val).replace('%','').replace(',','').strip())
    except Exception:
        return 0.0

df_equity['portfolio_ann'] = df_equity.iloc[:, 11].apply(safe_float)   # Column L (index 11)
df_equity['benchmark_ann'] = df_equity.iloc[:, 14].apply(safe_float)   # Column O (index 14)
df_equity['outperf'] = df_equity['portfolio_ann'] - df_equity['benchmark_ann']

# Calculate % allocation for ALL schemes (relative to global total)
df_equity['allocation_pct'] = df_equity['endmv'] / total_mv * 100
df_equity['weighted_alpha'] = df_equity['outperf'] * df_equity['endmv']

# Sort and get top 7 outperforming and top 1 underperforming
df_outperforming = df_equity[df_equity['outperf'] > 0].sort_values(by='outperf', ascending=False).head(7)
df_underperforming = df_equity[df_equity['outperf'] <= 0].sort_values(by='outperf').head(1)

def format_inr_cr(val):
    try:
        return f"{float(val)/1e7:,.2f}"
    except Exception:
        return ""

def format_pct(val):
    try:
        return f"{float(val):.2f}%"
    except Exception:
        return ""

def ensure_rows(table, n):
    while len(table.rows) < n:
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''
        table._tbl.append(new_row)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

# --- 2. Open PowerPoint, update tables, chart, and alpha text ---
slide = prs.slides[17]  # Slide 18 (0-based index)

# Find tables
tables = []
for shape in slide.shapes:
    if shape.has_table:
        tables.append(shape.table)
if len(tables) < 2:
    raise Exception("Expected two tables on slide 18!")

table_out = tables[0]
table_under = tables[1]

# Ensure enough rows in both tables
ensure_rows(table_out, len(df_outperforming)+2)
ensure_rows(table_under, len(df_underperforming)+2)

# --- 3. Fill Outperforming Table (top 7) ---
for i, (_, row) in enumerate(df_outperforming.iterrows(), start=1):
    set_cell(table_out.cell(i, 0), row[first_col], bold=False)
    set_cell(table_out.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_out.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_out.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for outperforming (including sum of outperformance)
set_cell(table_out.cell(len(df_outperforming)+1, 0), "Total", bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 1), format_inr_cr(df_outperforming['endmv'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 2), format_pct(df_outperforming['allocation_pct'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 3), format_pct(df_outperforming['outperf'].sum()), bold=True)

# --- 4. Fill Underperforming Table (top 1) ---
for i, (_, row) in enumerate(df_underperforming.iterrows(), start=1):
    set_cell(table_under.cell(i, 0), row[first_col], bold=False)
    set_cell(table_under.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_under.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_under.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for underperforming
set_cell(table_under.cell(len(df_underperforming)+1, 0), "Total", bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 1), format_inr_cr(df_underperforming['endmv'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 2), format_pct(df_underperforming['allocation_pct'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 3), format_pct(df_underperforming['outperf'].sum()), bold=True)

# --- 5. Update Pie Chart (by number of schemes, not value, rounded to 2 decimals) ---
num_outperf = (df_equity['outperf'] > 0).sum()
num_underperf = (df_equity['outperf'] <= 0).sum()
total_schemes = len(df_equity)

percent_schemes_outperf = round(num_outperf / total_schemes * 100, 2) if total_schemes else 0
percent_schemes_underperf = round(num_underperf / total_schemes * 100, 2) if total_schemes else 0

chart_data = CategoryChartData()
chart_data.categories = ['Outperforming', 'Underperforming']
chart_data.add_series('Schemes', (percent_schemes_outperf, percent_schemes_underperf))

for shape in slide.shapes:
    if hasattr(shape, "has_chart") and shape.has_chart:
        chart = shape.chart
        chart.replace_data(chart_data)
        break

# --- 6. Update the alpha text box (look for text containing "Equity Portfolio Alpha") ---
# Weighted average alpha: sum((outperf) * (endmv)) / total_mv
df_equity['weighted_alpha'] = df_equity['outperf'] * df_equity['endmv']
portfolio_alpha = df_equity['weighted_alpha'].sum() / total_mv if total_mv else 0

for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "Equity Portfolio Alpha" in run.text:
                    run.text = f"Equity Portfolio Alpha ~{portfolio_alpha:.2f}%"
                    break

prs.save("output_filled.pptx")
print("PowerPoint table, chart, and alpha text updated and saved as output_filled.pptx")

"""## Slide 22"""

# --- 1. Load Excel data ---
excel_path = "Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance SI", header=5)

# Normalize columns and first column for matching
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()
first_col = df.columns[0]
df[first_col] = df[first_col].astype(str).str.strip()

# Find start and stop indices for Debt section in first column
debt_idx = df[df[first_col].str.lower() == 'debt'].index
hybrid_idx = df[df[first_col].str.lower() == 'hybrid'].index

if len(debt_idx) == 0:
    raise Exception("No 'Debt' row found in first column.")
if len(hybrid_idx) == 0:
    stop_idx = len(df)
else:
    stop_idx = hybrid_idx[0]

start_idx = debt_idx[0]

# Slice the dataframe for Debt schemes only (exclusive of headers)
df_debt = df.loc[start_idx+1:stop_idx-1].copy()

# Skip unwanted scheme names
skip_names = [
    'debt mf', 'total : pms', 'total : aif', 'total : bonds / debentures',
    'total : mld', 'total : debt'
]
df_debt = df_debt[~df_debt[first_col].str.lower().isin(skip_names)]

# Convert endmv to numeric (handle commas, NaN, blanks)
df_debt['endmv'] = pd.to_numeric(
    df_debt['endmv'].astype(str).str.replace(',', '').replace('nan', ''), errors='coerce'
).fillna(0)

# Only schemes with End MV > 0
df_debt = df_debt[df_debt['endmv'] > 0]

# Calculate total End MV for allocation %
total_mv = df_debt['endmv'].sum()

# Calculate Portfolio and Benchmark Annualized Returns and Outperformance
def safe_float(val):
    try:
        return float(str(val).replace('%','').replace(',','').strip())
    except Exception:
        return 0.0

df_debt['portfolio_ann'] = df_debt.iloc[:, 11].apply(safe_float)   # Column L (index 11)
df_debt['benchmark_ann'] = df_debt.iloc[:, 14].apply(safe_float)   # Column O (index 14)
df_debt['outperf'] = df_debt['portfolio_ann'] - df_debt['benchmark_ann']

# Calculate % allocation for ALL schemes (relative to global total)
df_debt['allocation_pct'] = df_debt['endmv'] / total_mv * 100
df_debt['weighted_alpha'] = df_debt['outperf'] * df_debt['endmv']

# Sort and get top 7 outperforming and top 1 underperforming
df_outperforming = df_debt[df_debt['outperf'] > 0].sort_values(by='outperf', ascending=False).head(7)
df_underperforming = df_debt[df_debt['outperf'] <= 0].sort_values(by='outperf').head(1)

def format_inr_cr(val):
    try:
        return f"{float(val)/1e7:,.2f}"
    except Exception:
        return ""

def format_pct(val):
    try:
        return f"{float(val):.2f}%"
    except Exception:
        return ""

def ensure_rows(table, n):
    while len(table.rows) < n:
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''
        table._tbl.append(new_row)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

# --- 2. Open PowerPoint, update tables, chart, and alpha text ---
slide = prs.slides[21]  # Slide 22 (0-based index)

# Find tables
tables = []
for shape in slide.shapes:
    if shape.has_table:
        tables.append(shape.table)
if len(tables) < 2:
    raise Exception("Expected two tables on slide 22!")

table_out = tables[0]
table_under = tables[1]

# Ensure enough rows in both tables
ensure_rows(table_out, len(df_outperforming)+2)
ensure_rows(table_under, len(df_underperforming)+2)

# --- 3. Fill Outperforming Table (top 7) ---
for i, (_, row) in enumerate(df_outperforming.iterrows(), start=1):
    set_cell(table_out.cell(i, 0), row[first_col], bold=False)
    set_cell(table_out.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_out.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_out.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for outperforming (including sum of outperformance)
set_cell(table_out.cell(len(df_outperforming)+1, 0), "Total", bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 1), format_inr_cr(df_outperforming['endmv'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 2), format_pct(df_outperforming['allocation_pct'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 3), format_pct(df_outperforming['outperf'].sum()), bold=True)

# --- 4. Fill Underperforming Table (top 1) ---
for i, (_, row) in enumerate(df_underperforming.iterrows(), start=1):
    set_cell(table_under.cell(i, 0), row[first_col], bold=False)
    set_cell(table_under.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_under.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_under.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for underperforming
set_cell(table_under.cell(len(df_underperforming)+1, 0), "Total", bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 1), format_inr_cr(df_underperforming['endmv'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 2), format_pct(df_underperforming['allocation_pct'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 3), format_pct(df_underperforming['outperf'].sum()), bold=True)

# --- 5. Update Pie Chart (by number of schemes, not value, rounded to 2 decimals) ---
num_outperf = (df_debt['outperf'] > 0).sum()
num_underperf = (df_debt['outperf'] <= 0).sum()
total_schemes = len(df_debt)

percent_schemes_outperf = round(num_outperf / total_schemes * 100, 2) if total_schemes else 0
percent_schemes_underperf = round(num_underperf / total_schemes * 100, 2) if total_schemes else 0

chart_data = CategoryChartData()
chart_data.categories = ['Outperforming', 'Underperforming']
chart_data.add_series('Schemes', (percent_schemes_outperf, percent_schemes_underperf))

for shape in slide.shapes:
    if hasattr(shape, "has_chart") and shape.has_chart:
        chart = shape.chart
        chart.replace_data(chart_data)
        break

# --- 6. Update the alpha text box (look for text containing "Debt Portfolio Alpha") ---
df_debt['weighted_alpha'] = df_debt['outperf'] * df_debt['endmv']
portfolio_alpha = df_debt['weighted_alpha'].sum() / total_mv if total_mv else 0

for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "Debt Portfolio Alpha" in run.text:
                    run.text = f"Debt Portfolio Alpha ~{portfolio_alpha:.2f}%"
                    break

prs.save("output_filled.pptx")
print("PowerPoint table, chart, and alpha text updated and saved as output_filled.pptx")

"""## Slide 25"""

# --- 1. Load Excel data ---
excel_path = "Automation Reports.xlsx"
df = pd.read_excel(excel_path, sheet_name="Performance SI", header=5)

# Normalize columns and first column for matching
df.columns = df.columns.str.replace(' ', '', regex=False).str.lower()
first_col = df.columns[0]
df[first_col] = df[first_col].astype(str).str.strip()

# Find start and stop indices for Hybrid section in first column
hybrid_idx = df[df[first_col].str.lower() == 'hybrid'].index
alt_idx = df[df[first_col].str.lower() == 'alternative'].index

if len(hybrid_idx) == 0:
    raise Exception("No 'Hybrid' row found in first column.")
if len(alt_idx) == 0:
    stop_idx = len(df)
else:
    stop_idx = alt_idx[0]

start_idx = hybrid_idx[0]

# Slice the dataframe for Hybrid schemes only (exclusive of headers)
df_hybrid = df.loc[start_idx+1:stop_idx-1].copy()

# Skip unwanted scheme names
skip_names = [
    'hybrid mf', 'total : aif', 'total : hybrid'
]
df_hybrid = df_hybrid[~df_hybrid[first_col].str.lower().isin(skip_names)]

# Convert endmv to numeric (handle commas, NaN, blanks)
df_hybrid['endmv'] = pd.to_numeric(
    df_hybrid['endmv'].astype(str).str.replace(',', '').replace('nan', ''), errors='coerce'
).fillna(0)

# Only schemes with End MV > 0
df_hybrid = df_hybrid[df_hybrid['endmv'] > 0]

# Calculate total End MV for allocation %
total_mv = df_hybrid['endmv'].sum()

# Calculate Portfolio and Benchmark Annualized Returns and Outperformance
def safe_float(val):
    try:
        return float(str(val).replace('%','').replace(',','').strip())
    except Exception:
        return 0.0

df_hybrid['portfolio_ann'] = df_hybrid.iloc[:, 11].apply(safe_float)   # Column L (index 11)
df_hybrid['benchmark_ann'] = df_hybrid.iloc[:, 14].apply(safe_float)   # Column O (index 14)
df_hybrid['outperf'] = df_hybrid['portfolio_ann'] - df_hybrid['benchmark_ann']

# Calculate % allocation for ALL schemes (relative to global total)
df_hybrid['allocation_pct'] = df_hybrid['endmv'] / total_mv * 100
df_hybrid['weighted_alpha'] = df_hybrid['outperf'] * df_hybrid['endmv']

# Sort and get top 7 outperforming and top 1 underperforming
df_outperforming = df_hybrid[df_hybrid['outperf'] > 0].sort_values(by='outperf', ascending=False).head(7)
df_underperforming = df_hybrid[df_hybrid['outperf'] <= 0].sort_values(by='outperf').head(1)

def format_inr_cr(val):
    try:
        return f"{float(val)/1e7:,.2f}"
    except Exception:
        return ""

def format_pct(val):
    try:
        return f"{float(val):.2f}%"
    except Exception:
        return ""

def ensure_rows(table, n):
    while len(table.rows) < n:
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''
        table._tbl.append(new_row)

def set_cell(cell, text, bold=False):
    cell.text = str(text) if text is not None else ""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(9)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

# --- 2. Open PowerPoint, update tables, chart, and alpha text ---
slide = prs.slides[24]  # Slide 25 (0-based index)

# Find tables
tables = []
for shape in slide.shapes:
    if shape.has_table:
        tables.append(shape.table)
if len(tables) < 2:
    raise Exception("Expected two tables on slide 25!")

table_out = tables[0]
table_under = tables[1]

# Ensure enough rows in both tables
ensure_rows(table_out, len(df_outperforming)+2)
ensure_rows(table_under, len(df_underperforming)+2)

# --- 3. Fill Outperforming Table (top 7) ---
for i, (_, row) in enumerate(df_outperforming.iterrows(), start=1):
    set_cell(table_out.cell(i, 0), row[first_col], bold=False)
    set_cell(table_out.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_out.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_out.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for outperforming (including sum of outperformance)
set_cell(table_out.cell(len(df_outperforming)+1, 0), "Total", bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 1), format_inr_cr(df_outperforming['endmv'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 2), format_pct(df_outperforming['allocation_pct'].sum()), bold=True)
set_cell(table_out.cell(len(df_outperforming)+1, 3), format_pct(df_outperforming['outperf'].sum()), bold=True)

# --- 4. Fill Underperforming Table (top 1) ---
for i, (_, row) in enumerate(df_underperforming.iterrows(), start=1):
    set_cell(table_under.cell(i, 0), row[first_col], bold=False)
    set_cell(table_under.cell(i, 1), format_inr_cr(row['endmv']), bold=False)
    set_cell(table_under.cell(i, 2), format_pct(row['allocation_pct']), bold=False)
    set_cell(table_under.cell(i, 3), format_pct(row['outperf']), bold=False)

# Fill Total row for underperforming
set_cell(table_under.cell(len(df_underperforming)+1, 0), "Total", bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 1), format_inr_cr(df_underperforming['endmv'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 2), format_pct(df_underperforming['allocation_pct'].sum()), bold=True)
set_cell(table_under.cell(len(df_underperforming)+1, 3), format_pct(df_underperforming['outperf'].sum()), bold=True)

# --- 5. Update Pie Chart (by number of schemes, not value, rounded to 2 decimals) ---
num_outperf = (df_hybrid['outperf'] > 0).sum()
num_underperf = (df_hybrid['outperf'] <= 0).sum()
total_schemes = len(df_hybrid)

percent_schemes_outperf = round(num_outperf / total_schemes * 100, 2) if total_schemes else 0
percent_schemes_underperf = round(num_underperf / total_schemes * 100, 2) if total_schemes else 0

chart_data = CategoryChartData()
chart_data.categories = ['Outperforming', 'Underperforming']
chart_data.add_series('Schemes', (percent_schemes_outperf, percent_schemes_underperf))

for shape in slide.shapes:
    if hasattr(shape, "has_chart") and shape.has_chart:
        chart = shape.chart
        chart.replace_data(chart_data)
        break

# --- 6. Update the alpha text box (look for text containing "Hybrid Portfolio Alpha") ---
df_hybrid['weighted_alpha'] = df_hybrid['outperf'] * df_hybrid['endmv']
portfolio_alpha = df_hybrid['weighted_alpha'].sum() / total_mv if total_mv else 0

for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "Hybrid Portfolio Alpha" in run.text:
                    run.text = f"Hybrid Portfolio Alpha ~{portfolio_alpha:.2f}%"
                    break

prs.save("output_filled.pptx")
print("PowerPoint table, chart, and alpha text updated and saved as output_filled.pptx")

def get_table_info(ppt_path, slide_index):
    prs = Presentation(ppt_path)
    slide = prs.slides[slide_index]
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    if not table:
        print("No table found on this slide.")
        return
    num_rows = len(table.rows)
    num_cols = len(table.columns)
    content = []
    for row in table.rows:
        row_content = [cell.text for cell in row.cells]
        content.append(row_content)
    print("Number of rows:", num_rows)
    print("Number of columns:", num_cols)
    print("Table content:")
    for row in content:
        print(row)
    return num_rows, num_cols, content

ppt_path = "/content/Quarterly Review Revised Template_new (1).pptx"
slide_index = 13  # slide 14 (0-based index)
get_table_info(ppt_path, slide_index)