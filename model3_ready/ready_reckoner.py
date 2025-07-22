import os
import tempfile
import re
import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION

def make_filename_safe(name):
    return re.sub(r'[\\/*?:"<>|]', "_", name)

def set_cell_font_and_center(cell, text, font_name='Calibri', font_size=11, bold=False):
    cell.text = str(text) if text is not None else ''
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold

def set_cell_font_and_center_percent(cell, text, font_name='Calibri', font_size=11, bold=False):
    try:
        if isinstance(text, (float, int)):
            text = f"{text*100:.2f}%"
        elif isinstance(text, str):
            num = float(str(text).replace('%', '').replace(',', ''))
            text = f"{num*100:.2f}%"
    except:
        text = str(text) if text is not None else ''
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold

def fill_top_holdings_table_pms(table, fund_name, df_holding):
    fund_row = df_holding.iloc[1]
    fund_col_idx = None
    for i, val in enumerate(fund_row):
        if pd.notna(val) and str(val).strip().lower() == fund_name.strip().lower():
            fund_col_idx = i
            break
    if fund_col_idx is None:
        return
    stock_data = df_holding.iloc[17:27, [fund_col_idx, fund_col_idx+1]]
    for i, (_, row) in enumerate(stock_data.iterrows(), start=2):
        stock_name = row.iloc[0]
        stock_alloc = row.iloc[1]
        set_cell_font_and_center(table.cell(i, 0), stock_name)
        try:
            if pd.isna(stock_alloc) or stock_alloc == "":
                alloc_text = ""
            else:
                alloc_text = f"{float(stock_alloc)*100:.2f}%"
        except:
            alloc_text = str(stock_alloc) if stock_alloc is not None else ""
        set_cell_font_and_center(table.cell(i, 1), alloc_text)

def add_sector_allocation_chart_pms(slide, fund_name, df_holding):
    fund_row = df_holding.iloc[1]
    fund_col_idx = None
    for i, val in enumerate(fund_row):
        if pd.notna(val) and str(val).strip().lower() == fund_name.strip().lower():
            fund_col_idx = i
            break
    if fund_col_idx is None:
        return
    sector_data = df_holding.iloc[3:13, [fund_col_idx, fund_col_idx+1]].dropna(how='all')
    sectors = [str(x) for x in sector_data.iloc[:, 0]]
    allocations = []
    for val in sector_data.iloc[:, 1]:
        try:
            allocations.append(float(val)*100)
        except:
            allocations.append(0)
    chart_data = CategoryChartData()
    chart_data.categories = sectors
    chart_data.add_series('Sector Allocation', allocations)
    left = Inches(8)
    top = Inches(4.45)
    width = Inches(4.8)
    height = Inches(3)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart
    series = chart.series[0]
    for point in series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor(0x08, 0x47, 0x95)
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    value_axis = chart.value_axis
    value_axis.maximum_scale = 50.0
    value_axis.tick_labels.number_format = '0"%"'
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.bold = False
    value_axis.tick_labels.font.name = "Calibri"
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.bold = False
    category_axis.tick_labels.font.name = "Calibri"
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis.tick_labels.rotation = 90
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.number_format = '0.00'
    data_labels.font.bold = True
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.has_title = True
    chart.chart_title.text_frame.text = "Sectoral Allocation"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = "Calibri"

def fill_top_holdings_table(table, fund_name, df_holding):
    fund_row = df_holding.iloc[1]
    fund_col_idx = None
    for i, val in enumerate(fund_row):
        if pd.notna(val) and str(val).strip().lower() == fund_name.strip().lower():
            fund_col_idx = i
            break
    if fund_col_idx is None:
        return
    stock_data = df_holding.iloc[19:28, [fund_col_idx, fund_col_idx+1]]
    for i, (_, row) in enumerate(stock_data.iterrows(), start=2):
        stock_name = row.iloc[0]
        stock_alloc = row.iloc[1]
        set_cell_font_and_center(table.cell(i, 0), stock_name)
        try:
            if pd.isna(stock_alloc) or stock_alloc == "":
                alloc_text = ""
            else:
                alloc_text = f"{float(stock_alloc)*100:.2f}%"
        except:
            alloc_text = str(stock_alloc) if stock_alloc is not None else ""
        set_cell_font_and_center(table.cell(i, 1), alloc_text)

def add_sector_allocation_chart(slide, fund_name, df_holding):
    fund_row = df_holding.iloc[1]
    fund_col_idx = None
    for i, val in enumerate(fund_row):
        if pd.notna(val) and str(val).strip().lower() == fund_name.strip().lower():
            fund_col_idx = i
            break
    if fund_col_idx is None:
        return
    sector_data = df_holding.iloc[4:14, [fund_col_idx, fund_col_idx+1]].dropna(how='all')
    sectors = [str(x) for x in sector_data.iloc[:, 0]]
    allocations = []
    for val in sector_data.iloc[:, 1]:
        try:
            allocations.append(float(val)*100)
        except:
            allocations.append(0)
    chart_data = CategoryChartData()
    chart_data.categories = sectors
    chart_data.add_series('Sector Allocation', allocations)
    left = Inches(8)
    top = Inches(4.45)
    width = Inches(4.8)
    height = Inches(3)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart
    series = chart.series[0]
    for point in series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor(192, 79, 21)
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    value_axis = chart.value_axis
    value_axis.maximum_scale = 50.0
    value_axis.tick_labels.number_format = '0.00"%"'
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.bold = False
    value_axis.tick_labels.font.name = "Calibri"
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.bold = False
    category_axis.tick_labels.font.name = "Calibri"
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis.tick_labels.rotation = 90
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.number_format = '0.00"%"'
    data_labels.font.bold = True
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.has_title = True
    chart.chart_title.text_frame.text = "Sectoral Allocation"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = "Calibri"

def run_pms(excel_path, pms_template_path, n_pms, output_path):
    prs = Presentation(pms_template_path)
    df_pms = pd.read_excel(excel_path, sheet_name="PMS", header=2)
    df_perf = pd.read_excel(excel_path, sheet_name="Performance Formatted", header=0)
    df_bench = pd.read_excel(excel_path, sheet_name="Benchmarks", header=0)
    df_holding = pd.read_excel(excel_path, sheet_name="Holding", header=None)
    cy_cols = ["YTD", "CY 2024", "CY 2023", 2022, 2021, 2020]
    perf_cols = ["1 month", "3 months", "6 months", "1 Year", "3 Year", "5 Year"]
    fund_list = df_pms["Fund Name"].dropna().unique().tolist()

    for idx, fund_name in enumerate(fund_list[:n_pms]):
        slide = prs.slides[idx]
        tables = [shape for shape in slide.shapes if shape.has_table]
        matching_rows = df_pms[df_pms["Fund Name"].astype(str).str.strip().str.lower() == fund_name.lower()]
        if matching_rows.empty: continue
        fund_row = matching_rows.iloc[0]
        # Table 0: CY Returns
        cy_table = tables[0].table
        perf_row = df_perf[df_perf["Fund Name"].astype(str).str.strip().str.lower() == fund_name.lower()]
        fund_cy_data = [perf_row.iloc[0][col] if not perf_row.empty and col in perf_row.columns else "" for col in cy_cols]
        bench_row = df_bench[df_bench["Fund Name"].astype(str).str.strip().str.lower() == fund_name.lower()]
        if not bench_row.empty:
            benchmark_name = bench_row.iloc[0]["Benchmark"]
            bench_cy_data = [bench_row.iloc[0][col] if col in bench_row.columns else "" for col in cy_cols]
        else:
            benchmark_name = "Benchmark"
            bench_cy_data = [""] * len(cy_cols)
        set_cell_font_and_center_percent(cy_table.cell(2, 0), fund_name, bold=True)
        set_cell_font_and_center_percent(cy_table.cell(3, 0), benchmark_name, bold=True)
        for col_idx, val in enumerate(fund_cy_data, start=1):
            set_cell_font_and_center_percent(cy_table.cell(2, col_idx), val)
        for col_idx, val in enumerate(bench_cy_data, start=1):
            set_cell_font_and_center_percent(cy_table.cell(3, col_idx), val)
        # Table 1: Periodical Performance
        perf_table = tables[1].table
        fund_perf_data = [perf_row.iloc[0][col] if not perf_row.empty and col in perf_row.columns else "" for col in perf_cols]
        bench_perf_data = [bench_row.iloc[0][col] if not bench_row.empty and col in bench_row.columns else "" for col in perf_cols]
        set_cell_font_and_center_percent(perf_table.cell(2, 0), fund_name, bold=True)
        set_cell_font_and_center_percent(perf_table.cell(3, 0), benchmark_name, bold=True)
        for col_idx, val in enumerate(fund_perf_data, start=1):
            set_cell_font_and_center_percent(perf_table.cell(2, col_idx), val)
        for col_idx, val in enumerate(bench_perf_data, start=1):
            set_cell_font_and_center_percent(perf_table.cell(3, col_idx), val)
        # Fund Details Table (Table 5)
        set_cell_font_and_center(tables[5].table.cell(1, 1), fund_row.get("Investment Strategy",''))
        set_cell_font_and_center(tables[5].table.cell(2, 1), fund_row.get("Category",''))
        set_cell_font_and_center(tables[5].table.cell(3, 1), fund_row.get("AUM (Crs)",''))
        set_cell_font_and_center(tables[5].table.cell(4, 1), fund_row.get("Portfolio Style",''))
        set_cell_font_and_center(tables[5].table.cell(5, 1), fund_row.get("Expense Ratio- Regular",''))
        set_cell_font_and_center(tables[5].table.cell(6, 1), fund_row.get("Exit Load",''))
        set_cell_font_and_center(tables[5].table.cell(7, 1), fund_row.get("Fund Manager",''))
        # Risk Ratios Table (Table 2)
        set_cell_font_and_center(tables[2].table.cell(1, 1), fund_row.get("Standard Deviation (%)",''))
        set_cell_font_and_center(tables[2].table.cell(2, 1), fund_row.get("Sharpe Ratio",''))
        set_cell_font_and_center(tables[2].table.cell(3, 1), fund_row.get("Beta",''))
        # Portfolio Table (Table 3)
        set_cell_font_and_center(tables[3].table.cell(1, 1), fund_row.get("No. of stocks in portfolio",''))
        set_cell_font_and_center(tables[3].table.cell(2, 1), fund_row.get("Top 10 stocks (%)",''))
        set_cell_font_and_center(tables[3].table.cell(3, 1), fund_row.get("Top 5 stocks (%)",''))
        set_cell_font_and_center(tables[3].table.cell(4, 1), fund_row.get("Top 3 sectors (%)",''))
        # Valuations Table (Table 4)
        set_cell_font_and_center(tables[4].table.cell(1, 1), fund_row.get("P/B Ratio",''))
        set_cell_font_and_center(tables[4].table.cell(2, 1), fund_row.get("P/E Ratio",''))
        # MCap Allocation Table (Table 7)
        set_cell_font_and_center(tables[7].table.cell(1, 1), fund_row.get("Large (%)",''))
        set_cell_font_and_center(tables[7].table.cell(2, 1), fund_row.get("Mid (%)",''))
        set_cell_font_and_center(tables[7].table.cell(3, 1), fund_row.get("Small (%)",''))
        # Fill Top Holdings Table (Table 6)
        fill_top_holdings_table_pms(tables[6].table, fund_name, df_holding)
        # Sector allocation chart
        add_sector_allocation_chart_pms(slide, fund_name, df_holding)
        # Fund name heading
        left = Inches(0.5)
        top = Inches(0.18)
        width = Inches(7.5)
        height = Inches(0.8)
        heading_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = heading_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = fund_name
        run.font.name = "Roboto"
        run.font.size = Pt(32)
        p.alignment = PP_ALIGN.LEFT
        run.font.color.rgb = RGBColor(8, 71, 149)
    prs.save(output_path)

def run_hybrid(excel_path, hybrid_template_path, n_hybrid, output_path):
    prs = Presentation(hybrid_template_path)
    df_hybrid = pd.read_excel(excel_path, sheet_name="Hybrid MF", header=2)
    df_hybrid_hold = pd.read_excel(excel_path, sheet_name="Hybrid MF Holding", header=None)
    df_perf = pd.read_excel(excel_path, sheet_name="Performance Formatted", header=0)
    df_bench = pd.read_excel(excel_path, sheet_name="Benchmarks", header=0)
    cy_cols = ["YTD", "CY 2024", "CY 2023", 2022, 2021, 2020]
    perf_cols = ["1 month", "3 months", "6 months", "1 Year", "3 Year", "5 Year"]
    hybrid_col = None
    for col in df_hybrid.columns:
        if 'hybrid' in str(col).lower():
            hybrid_col = col
            break
    if hybrid_col is None:
        raise Exception("Hybrid column not found in Hybrid MF sheet")
    hybrid_fund_list = df_hybrid[hybrid_col].dropna().unique().tolist()
    slide_idx_offset = 0
    for idx, fund_name in enumerate(hybrid_fund_list[:n_hybrid]):
        slide = prs.slides[slide_idx_offset + idx]
        tables = [shape for shape in slide.shapes if shape.has_table]
        # Table 0: CY Returns
        cy_table = tables[0].table
        perf_row = df_perf[df_perf["Fund Name"].astype(str).str.strip().str.lower() == fund_name.lower()]
        fund_cy_data = [perf_row.iloc[0][col] if not perf_row.empty and col in perf_row.columns else "" for col in cy_cols]
        bench_row = df_bench[df_bench["Fund Name"].astype(str).str.strip().str.lower() == fund_name.lower()]
        if not bench_row.empty:
            benchmark_name = bench_row.iloc[0]["Benchmark"]
            bench_cy_data = [bench_row.iloc[0][col] if col in bench_row.columns else "" for col in cy_cols]
        else:
            benchmark_name = "Benchmark"
            bench_cy_data = [""] * len(cy_cols)
        set_cell_font_and_center_percent(cy_table.cell(2, 0), fund_name, bold=True)
        set_cell_font_and_center_percent(cy_table.cell(3, 0), benchmark_name, bold=True)
        for col_idx, val in enumerate(fund_cy_data, start=1):
            set_cell_font_and_center_percent(cy_table.cell(2, col_idx), val)
        for col_idx, val in enumerate(bench_cy_data, start=1):
            set_cell_font_and_center_percent(cy_table.cell(3, col_idx), val)
        # Table 1: Periodical Performance
        perf_table = tables[1].table
        fund_perf_data = [perf_row.iloc[0][col] if not perf_row.empty and col in perf_row.columns else "" for col in perf_cols]
        bench_perf_data = [bench_row.iloc[0][col] if not bench_row.empty and col in bench_row.columns else "" for col in perf_cols]
        set_cell_font_and_center_percent(perf_table.cell(2, 0), fund_name, bold=True)
        set_cell_font_and_center_percent(perf_table.cell(3, 0), benchmark_name, bold=True)
        for col_idx, val in enumerate(fund_perf_data, start=1):
            set_cell_font_and_center_percent(perf_table.cell(2, col_idx), val)
        for col_idx, val in enumerate(bench_perf_data, start=1):
            set_cell_font_and_center_percent(perf_table.cell(3, col_idx), val)
        # Asset Class Table (Table 2)
        asset_class_table = tables[2].table
        v = lambda k: df_hybrid.loc[df_hybrid[hybrid_col] == fund_name, k].values[0] if len(df_hybrid.loc[df_hybrid[hybrid_col] == fund_name, k])>0 else ''
        set_cell_font_and_center(asset_class_table.cell(1, 1), v("Equity"))
        set_cell_font_and_center(asset_class_table.cell(2, 1), v("Debt"))
        set_cell_font_and_center(asset_class_table.cell(3, 1), v("Others"))
        # Portfolio Table (Table 3)
        portfolio_table = tables[3].table
        set_cell_font_and_center(portfolio_table.cell(1, 1), v("Gross YTM"))
        set_cell_font_and_center(portfolio_table.cell(2, 1), v("MoD"))
        set_cell_font_and_center(portfolio_table.cell(3, 1), v("Avg Maturity"))
        set_cell_font_and_center(portfolio_table.cell(4, 1), v("SOV"))
        set_cell_font_and_center(portfolio_table.cell(5, 1), v("AAA"))
        set_cell_font_and_center(portfolio_table.cell(6, 1), v("AA"))
        set_cell_font_and_center(portfolio_table.cell(7, 1), v("A and below"))
        set_cell_font_and_center(portfolio_table.cell(8, 1), v("Unrated"))
        set_cell_font_and_center(portfolio_table.cell(9, 1), v("Cash"))
        # Valuations Table (Table 4)
        valuations_table = tables[4].table
        set_cell_font_and_center(valuations_table.cell(1, 1), v("P/B"))
        set_cell_font_and_center(valuations_table.cell(2, 1), v("P/E"))
        # Fund Details Table (Table 5)
        fund_details_table = tables[5].table
        set_cell_font_and_center(fund_details_table.cell(1, 1), v("Investment Strategy"))
        set_cell_font_and_center(fund_details_table.cell(2, 1), v("Category"))
        set_cell_font_and_center(fund_details_table.cell(3, 1), v("AUM (Crs)"))
        set_cell_font_and_center(fund_details_table.cell(4, 1), v("Fund Strategy"))
        set_cell_font_and_center(fund_details_table.cell(5, 1), v("Expense Ratio (Direct)"))
        set_cell_font_and_center(fund_details_table.cell(6, 1), v("Expense Ratio (Regular)"))
        set_cell_font_and_center(fund_details_table.cell(7, 1), v("Exit Load"))
        set_cell_font_and_center(fund_details_table.cell(8, 1), v("Fund Manager"))
        # Top Holdings Table (Table 6)
        fill_top_holdings_table(tables[6].table, fund_name, df_hybrid_hold)
        # MCap Allocation Table (Table 7)
        set_cell_font_and_center(tables[7].table.cell(1, 1), v("Large Cap"))
        set_cell_font_and_center(tables[7].table.cell(2, 1), v("Mid Cap"))
        set_cell_font_and_center(tables[7].table.cell(3, 1), v("Small Cap"))
        # Sector allocation chart
        add_sector_allocation_chart(slide, fund_name, df_hybrid_hold)
        # Fund name heading (Orange)
        left = Inches(0.5)
        top = Inches(0.18)
        width = Inches(7.5)
        height = Inches(0.8)
        heading_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = heading_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = fund_name
        run.font.name = "Roboto"
        run.font.size = Pt(32)
        p.alignment = PP_ALIGN.LEFT
        run.font.color.rgb = RGBColor(192, 79, 21)
    prs.save(output_path)
